"""PPTX extractor (python-pptx). One segment per slide."""
from __future__ import annotations

from pathlib import Path

from base import ExtractResult, ExtractorUnavailable, Segment


def _load():
    try:
        import pptx
        return pptx
    except ImportError as e:
        raise ExtractorUnavailable(
            "python-pptx not installed (pip install -r extractors/requirements.txt)") from e


def extract(path, granularity="structural"):
    pptx = _load()
    prs = pptx.Presentation(str(path))
    source = Path(path).name
    segs = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        title = (texts[0].splitlines()[0][:80] if texts else f"Slide {i}")
        segs.append(Segment(f"slide{i}", title, "\n\n".join(texts),
                            {"source": source, "slide": i}))
    return ExtractResult(Path(path).stem, "pptx", segs)
